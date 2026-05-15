#!/usr/bin/env python3
"""Build the final RQ47 temporal awareness presentation.

The deck is intentionally generated from existing artifacts. It does not rerun
model experiments or mutate any existing research code.
"""

from __future__ import annotations

import json
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable

import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
OUT_DIR = ROOT / "presentations"
PPTX_PATH = OUT_DIR / "rq47_temporal_awareness_final.pptx"
NOTES_PATH = OUT_DIR / "rq47_temporal_awareness_final_notes.md"

TABLE_DIR = ROOT / "results" / "tables" / "rq47"
PROBE_FIG_DIR = ROOT / "results" / "figures" / "probe_validation_multimethod"
PROBE_ALL_FIG_DIR = ROOT / "results" / "figures" / "probe_validation_all_methods"
RQ47_FIG_DIR = ROOT / "results" / "figures" / "rq47"

SLIDE_W = 13.333
SLIDE_H = 7.5

NAVY = RGBColor(23, 38, 63)
INK = RGBColor(31, 41, 55)
MUTED = RGBColor(99, 110, 123)
LIGHT = RGBColor(244, 247, 251)
BLUE = RGBColor(44, 104, 171)
TEAL = RGBColor(37, 139, 126)
ORANGE = RGBColor(204, 93, 50)
RED = RGBColor(172, 55, 55)
GREEN = RGBColor(65, 133, 88)
WHITE = RGBColor(255, 255, 255)


@dataclass
class SlideNotes:
    title: str
    bullets: list[str]


def required_path(path: Path) -> Path:
    if not path.exists():
        raise FileNotFoundError(f"Missing required presentation artifact: {path.relative_to(ROOT)}")
    return path


def read_csv(name: str) -> pd.DataFrame:
    return pd.read_csv(required_path(TABLE_DIR / name))


def add_text(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    size: int = 20,
    bold: bool = False,
    color: RGBColor = INK,
    align=PP_ALIGN.LEFT,
):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    add_text(slide, title, 0.55, 0.28, 11.7, 0.45, size=22, bold=True, color=NAVY)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(0.86), Inches(12.2), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()
    if subtitle:
        add_text(slide, subtitle, 0.55, 0.93, 11.8, 0.35, size=11, color=MUTED)


def add_bullets(
    slide,
    bullets: Iterable[str],
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    size: int = 17,
    color: RGBColor = INK,
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(7)


def add_card(slide, left: float, top: float, width: float, height: float, fill: RGBColor = WHITE):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = RGBColor(219, 226, 235)
    shape.line.width = Pt(1)
    return shape


def add_metric_card(slide, label: str, value: str, left: float, top: float, width: float, color: RGBColor) -> None:
    add_card(slide, left, top, width, 0.9)
    add_text(slide, value, left + 0.15, top + 0.12, width - 0.3, 0.33, size=21, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(slide, label, left + 0.15, top + 0.52, width - 0.3, 0.25, size=9, color=MUTED, align=PP_ALIGN.CENTER)


def add_explanation(
    slide,
    heading: str,
    body: str,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    color: RGBColor = BLUE,
) -> None:
    add_card(slide, left, top, width, height, fill=RGBColor(248, 250, 253))
    add_text(slide, heading, left + 0.18, top + 0.16, width - 0.36, 0.28, size=12, bold=True, color=color)
    add_text(slide, body, left + 0.18, top + 0.52, width - 0.36, height - 0.65, size=11, color=INK)


def add_image(slide, path: Path, left: float, top: float, width: float, height: float) -> None:
    required_path(path)
    with Image.open(path) as image:
        img_w, img_h = image.size
    img_ratio = img_w / img_h
    box_ratio = width / height
    if img_ratio >= box_ratio:
        draw_w = width
        draw_h = width / img_ratio
    else:
        draw_h = height
        draw_w = height * img_ratio
    draw_left = left + (width - draw_w) / 2
    draw_top = top + (height - draw_h) / 2
    slide.shapes.add_picture(str(path), Inches(draw_left), Inches(draw_top), Inches(draw_w), Inches(draw_h))


def add_simple_table(
    slide,
    df: pd.DataFrame,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    font_size: int = 9,
) -> None:
    rows, cols = df.shape[0] + 1, df.shape[1]
    table = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height)).table
    for c, name in enumerate(df.columns):
        cell = table.cell(0, c)
        cell.text = str(name)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(font_size)
            p.font.color.rgb = WHITE
    for r, row in enumerate(df.itertuples(index=False), start=1):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 else LIGHT
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.color.rgb = INK
                p.alignment = PP_ALIGN.CENTER


def add_flow_box(slide, label: str, left: float, top: float, width: float, color: RGBColor) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.85))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    add_text(slide, label, left + 0.1, top + 0.22, width - 0.2, 0.35, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def add_arrow(slide, left: float, top: float) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(left), Inches(top), Inches(0.45), Inches(0.28))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(169, 180, 194)
    shape.line.fill.background()


def format_pct(value: float) -> str:
    return f"{value:.0%}"


def load_summary_bits() -> dict[str, object]:
    probe = read_csv("probe_best_layers.csv")
    steering = read_csv("steering_dose_response_summary.csv")
    activation = read_csv("activation_patching_summary.csv")
    attribution = read_csv("attribution_patching_summary.csv")
    ablation = read_csv("ablation_summary.csv")
    oversight = read_csv("temporal_oversight_summary.csv")
    drift = read_csv("temporal_oversight_top_pre_event_drifts.csv")

    validation_rows = []
    for method in ["lr", "dmm", "attn"]:
        for path in sorted((ROOT / "results" / "probe_validation" / method).glob(f"*_probe_validation_{method}.json")):
            data = json.loads(path.read_text())
            summary = data.get("summary", {})
            model = data.get("metadata", {}).get("model", {})
            if not summary or not model:
                continue
            validation_rows.append(
                {
                    "method": method,
                    "model_tag": model.get("tag", path.name.split("_probe_validation_")[0]),
                    "model_name": model.get("name", ""),
                    "best_layer": int(summary["best_semantic_layer"]),
                    "best_explicit_acc": float(summary["best_explicit_accuracy"]),
                    "mean_explicit_acc": float(summary["mean_explicit_accuracy"]),
                    "n_semantic_layers": len(summary["semantic_layers"]),
                    "validation_passed": bool(summary["validation_passed"]),
                }
            )
    validation = pd.DataFrame(validation_rows)

    probe_pivot = probe.pivot(index="method", columns="model_tag", values="best_test_acc")
    hidden_probe_best = probe[probe["method"].isin(["lr", "dmm"])]["best_test_acc"].max()
    attn_min = probe[probe["method"] == "attn"]["best_test_acc"].min()
    attn_max = probe[probe["method"] == "attn"]["best_test_acc"].max()

    best_activation = (
        activation[(activation["layer_selection"] == "best") & (activation["component"] == "resid")]
        .sort_values("median_effect", ascending=False)
        .head(5)
    )
    best_attribution = (
        attribution[(attribution["layer_selection"] == "best") & (attribution["component"] == "resid")]
        .sort_values("median_effect", ascending=False)
        .head(5)
    )
    best_steering = (
        steering[(steering["layer_selection"] == "best") & (steering["component"] == "resid")]
        .sort_values("plus3_minus_zero", ascending=False)
        .head(4)
    )
    best_ablation = (
        ablation[(ablation["layer_selection"] == "best") & (ablation["component"] == "resid")]
        .sort_values("median_drop", ascending=False)
        .head(4)
    )
    top_drift = drift.head(6).copy()
    top_drift["delta_before_event"] = top_drift["delta_before_event"].map(lambda v: f"{v:.2f}")
    top_drift["first_event_step"] = top_drift["first_event_step"].map(lambda v: str(int(v)) if pd.notna(v) else "-")

    return {
        "probe": probe,
        "validation": validation,
        "probe_pivot": probe_pivot,
        "hidden_probe_best": hidden_probe_best,
        "attn_min": attn_min,
        "attn_max": attn_max,
        "best_activation": best_activation,
        "best_attribution": best_attribution,
        "best_steering": best_steering,
        "best_ablation": best_ablation,
        "oversight": oversight,
        "top_drift": top_drift,
    }


def build_deck() -> tuple[int, int, list[SlideNotes]]:
    bits = load_summary_bits()
    validation = bits["validation"]
    if validation.empty:
        raise ValueError("No method-scoped probe validation JSON files found under results/probe_validation/{lr,dmm,attn}")

    best_overall = validation.loc[validation["best_explicit_acc"].idxmax()]
    method_means = validation.groupby("method")["mean_explicit_acc"].mean().sort_values(ascending=False)
    all_pass = bool(validation["validation_passed"].all())

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank = prs.slide_layouts[6]
    notes: list[SlideNotes] = []

    def new_slide(title: str, subtitle: str | None = None):
        slide = prs.slides.add_slide(blank)
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor(250, 252, 255)
        add_title(slide, title, subtitle)
        return slide

    # 1
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY
    add_text(slide, "Temporal Horizon Representations", 0.75, 1.45, 11.8, 0.75, size=34, bold=True, color=WHITE)
    add_text(slide, "as causal and online oversight signals", 0.78, 2.23, 11.4, 0.55, size=24, color=RGBColor(210, 225, 245))
    add_text(slide, "RQ47 final technical presentation", 0.8, 4.95, 5.6, 0.35, size=14, color=RGBColor(194, 205, 218))
    add_metric_card(slide, "models", "4", 7.0, 4.65, 1.4, BLUE)
    add_metric_card(slide, "probe methods", "3", 8.55, 4.65, 1.6, TEAL)
    add_metric_card(slide, "causal tests", "4", 10.3, 4.65, 1.55, ORANGE)
    notes.append(SlideNotes("Title", ["Set the story: temporal horizon starts as a decodable representation, then becomes a causal and online monitoring target."]))

    # 2
    slide = new_slide("The Whole Story", "A coherent path from representation evidence to oversight evidence.")
    stages = [
        ("1. Decode", "Can a probe read short-term vs long-term horizon?", BLUE),
        ("2. Validate", "Does the signal transfer to explicit temporal examples?", TEAL),
        ("3. Intervene", "Does changing the representation move model logits?", ORANGE),
        ("4. Monitor", "Does the signal move while behavior unfolds?", GREEN),
    ]
    for idx, (head, body, color) in enumerate(stages):
        left = 0.85 + idx * 3.05
        add_card(slide, left, 1.65, 2.65, 3.2, fill=WHITE)
        add_text(slide, head, left + 0.2, 1.95, 2.25, 0.3, size=16, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(slide, body, left + 0.25, 2.55, 2.15, 1.1, size=13, color=INK, align=PP_ALIGN.CENTER)
    add_text(slide, "The claim grows only when each stage supports the previous one.", 1.2, 5.55, 10.9, 0.45, size=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    notes.append(SlideNotes("The Whole Story", ["This gives the audience a mental map before figures arrive. Each section answers one stronger question than the last."]))

    # 3
    slide = new_slide("Experimental Map", "The artifacts line up with the research progression.")
    xs = [0.65, 3.15, 5.65, 8.15, 10.65]
    labels = ["Train probes", "Validate probes", "Select layers", "Intervene", "Monitor tokens"]
    colors = [BLUE, TEAL, RGBColor(95, 99, 179), ORANGE, GREEN]
    for idx, (x, label, color) in enumerate(zip(xs, labels, colors)):
        add_flow_box(slide, label, x, 1.8, 1.85, color)
        if idx < len(xs) - 1:
            add_arrow(slide, x + 1.92, 2.08)
    add_bullets(
        slide,
        [
            "Probe training/validation creates layer-wise evidence and saved probes.",
            "The notebooks 03.1 and 03.2 turn the method/model results into validation figures.",
            "RQ47 intervention and trajectory scripts test whether the probe signal matters and moves online.",
        ],
        1.0,
        3.45,
        11.0,
        1.65,
        size=17,
    )
    notes.append(SlideNotes("Experimental Map", ["Point out that the deck uses the same artifact chain as the codebase: scripts produce results, notebooks summarize, deck tells the story."]))

    # 4
    slide = new_slide("Dataset And Probe Setup", "Implicit temporal CAA pairs train the probe; explicit temporal examples test semantic transfer.")
    for idx, (title, color, bullets) in enumerate(
        [
            ("Training", BLUE, ["Implicit AB-randomized temporal CAA examples", "Immediate and long-term options become balanced labels", "Pair-aware split avoids separating related examples"]),
            ("Validation", TEAL, ["Reload saved implicit-trained probes", "Score explicit temporal examples layer by layer", "Choose best semantic layer by explicit accuracy"]),
            ("Why It Matters", ORANGE, ["Success means the probe learned temporal horizon, not only one dataset's wording", "Layer-wise results tell causal experiments where to intervene"]),
        ]
    ):
        left = 0.75 + idx * 4.15
        add_card(slide, left, 1.45, 3.65, 4.25)
        add_text(slide, title, left + 0.22, 1.75, 3.1, 0.35, size=18, bold=True, color=color)
        add_bullets(slide, bullets, left + 0.25, 2.3, 3.1, 2.55, size=13)
    notes.append(SlideNotes("Dataset And Probe Setup", ["Explain the key intuition: implicit training plus explicit validation is a semantic generalization test."]))

    # 5
    slide = new_slide("Probe Methods", "Three probes ask whether temporal horizon is visible in hidden states or attention summaries.")
    methods = pd.DataFrame(
        [
            ["LR", "Hidden state", "Learned boundary; strongest flexible linear baseline"],
            ["DMM", "Hidden state", "Simple temporal direction; useful as interpretable baseline"],
            ["AttnProbe", "Attention summaries", "Tests whether attention distribution statistics carry signal"],
        ],
        columns=["Probe", "Feature source", "Interpretation"],
    )
    add_simple_table(slide, methods, 0.8, 1.35, 11.7, 1.35, font_size=12)
    add_metric_card(slide, "best hidden-state probe", format_pct(bits["hidden_probe_best"]), 0.95, 3.25, 2.75, BLUE)
    add_metric_card(slide, "AttnProbe range", f"{format_pct(bits['attn_min'])}-{format_pct(bits['attn_max'])}", 4.05, 3.25, 2.75, ORANGE)
    add_explanation(
        slide,
        "Meaning",
        "LR and DMM both finding signal is important: one is a learned classifier, the other is a simple average-difference direction. Agreement makes the representation less likely to be an artifact of one probe family.",
        7.2,
        3.1,
        4.85,
        1.65,
        color=TEAL,
    )
    add_image(slide, PROBE_FIG_DIR / "lr_minus_dmm_delta.png", 4.1, 5.0, 4.8, 1.25)
    notes.append(SlideNotes("Probe Methods", ["Make LR/DMM agreement intuitive: a strong learned classifier and a simple mean direction often point to the same underlying separation."]))

    # 6
    slide = new_slide("Finding 1: The Signal Generalizes", "Best explicit validation accuracy across all model/method cases.")
    add_image(slide, PROBE_FIG_DIR / "best_explicit_accuracy_heatmap.png", 0.75, 1.2, 7.4, 4.45)
    add_explanation(
        slide,
        "What The Heatmap Means",
        f"Each cell is the best explicit validation accuracy for one probe/model pair. The best case is {best_overall['method'].upper()} on {best_overall['model_tag']} at layer {int(best_overall['best_layer'])}, with {best_overall['best_explicit_acc']:.0%} explicit accuracy.",
        8.45,
        1.35,
        3.85,
        1.8,
        color=BLUE,
    )
    add_explanation(
        slide,
        "Why It Matters",
        "The probes were trained on implicit examples. High explicit accuracy means the representation transfers to a different temporal framing, so the probe is not merely memorizing superficial wording.",
        8.45,
        3.45,
        3.85,
        1.8,
        color=TEAL,
    )
    add_text(slide, f"Validation threshold passed by all cases: {'yes' if all_pass else 'no'}", 8.55, 5.7, 3.55, 0.35, size=14, bold=True, color=GREEN if all_pass else RED)
    notes.append(SlideNotes("Finding 1: The Signal Generalizes", ["Tell the audience how to read cells, then why cross-dataset transfer is the first big result."]))

    # 7
    slide = new_slide("Finding 1 Detail: Method Comparison", "The earlier all-method notebook shows which probe wins per model.")
    add_image(slide, PROBE_ALL_FIG_DIR / "best_explicit_by_method_bars.png", 0.75, 1.18, 7.8, 4.35)
    add_explanation(
        slide,
        "What The Bars Mean",
        "For each model, the bars compare the best explicit accuracy reached by LR, DMM, and AttnProbe. The dashed line marks the 70% semantic-validation threshold.",
        8.8,
        1.45,
        3.45,
        1.65,
        color=BLUE,
    )
    add_explanation(
        slide,
        "Why It Matters",
        f"The highest average explicit accuracy by method is {method_means.index[0].upper()} at {method_means.iloc[0]:.0%}. Hidden-state probes usually dominate, while attention summaries still carry weaker but real signal.",
        8.8,
        3.55,
        3.45,
        1.8,
        color=ORANGE,
    )
    notes.append(SlideNotes("Finding 1 Detail: Method Comparison", ["This slide translates the heatmap into a model-by-model comparison."]))

    # 8
    slide = new_slide("Finding 2: The Signal Has A Depth Profile", "Normalized depth makes architectures with different layer counts comparable.")
    add_image(slide, PROBE_ALL_FIG_DIR / "normalized_depth_heatmaps.png", 0.65, 1.12, 11.8, 4.6)
    add_explanation(
        slide,
        "What This Means",
        "Each row is a model, each column is a normalized depth bin, and each panel is a probe method. Brighter cells mean stronger explicit temporal-horizon decoding at that relative depth.",
        0.95,
        5.88,
        5.6,
        0.85,
        color=BLUE,
    )
    add_explanation(
        slide,
        "Why It Matters",
        "This answers where to look. The causal experiments use probe-selected layers rather than arbitrary layer choices, so the validation analysis directly informs intervention design.",
        6.85,
        5.88,
        5.25,
        0.85,
        color=TEAL,
    )
    notes.append(SlideNotes("Finding 2: The Signal Has A Depth Profile", ["This replaces the previously cluttered layer-dynamics slide with one large readable figure and explanations below."]))

    # 9
    slide = new_slide("Finding 2 Detail: Semantic Layer Coverage", "Some models have broad temporal-horizon bands; others have narrower windows.")
    add_image(slide, PROBE_FIG_DIR / "semantic_layer_coverage_heatmap.png", 0.85, 1.25, 7.2, 3.75)
    add_explanation(
        slide,
        "What The Counts Mean",
        "A higher count means more layers clear the 70% explicit validation threshold. Broad coverage means the representation is accessible across a larger depth range.",
        8.5,
        1.4,
        3.65,
        1.6,
        color=BLUE,
    )
    add_explanation(
        slide,
        "Why It Matters",
        "If only one layer works, causal results may be fragile. If many layers work, the temporal signal is more distributed and layer selection is less brittle.",
        8.5,
        3.45,
        3.65,
        1.55,
        color=GREEN,
    )
    notes.append(SlideNotes("Finding 2 Detail: Semantic Layer Coverage", ["Explain coverage as robustness of where the signal can be read."]))

    # 10
    slide = new_slide("Finding 3: Explicit And Implicit Curves Track Layer Behavior", "A single-case view shows what the layer curves are doing.")
    add_image(slide, PROBE_ALL_FIG_DIR / "explicit_vs_implicit_lr_Qwen3-4B.png", 0.65, 1.18, 8.0, 4.15)
    add_explanation(
        slide,
        "What To Look For",
        "The left panel compares implicit training/test behavior with explicit validation by layer. The right panel shows the explicit-minus-implicit gap and colors layers by semantic strength.",
        8.95,
        1.35,
        3.25,
        1.75,
        color=BLUE,
    )
    add_explanation(
        slide,
        "Why It Matters",
        "Layer curves reveal whether validation success is isolated, broad, or shifted relative to implicit training. That prevents a single best-layer number from hiding the shape of the result.",
        8.95,
        3.55,
        3.25,
        1.75,
        color=TEAL,
    )
    notes.append(SlideNotes("Finding 3: Explicit And Implicit Curves Track Layer Behavior", ["Use Qwen LR as a readable example of the earlier notebooks' per-case diagnostics."]))

    # 11
    slide = new_slide("Finding 4: Cross-Validation Supports The Probe Signal", "CV curves add uncertainty around the implicit-training probe performance.")
    add_image(slide, PROBE_ALL_FIG_DIR / "cv_analysis_lr.png", 0.65, 1.15, 8.0, 4.05)
    add_explanation(
        slide,
        "What The Error Bars Mean",
        "The LR CV figure reports pair-aware cross-validation on implicit examples, alongside implicit and explicit layer accuracy. Stable high CV reduces concern that a layer is just a split accident.",
        8.95,
        1.35,
        3.25,
        1.85,
        color=BLUE,
    )
    add_explanation(
        slide,
        "Why It Matters",
        "The deck's later causal claims depend on choosing credible layers. CV diagnostics help justify those layer choices before interventions are run.",
        8.95,
        3.65,
        3.25,
        1.55,
        color=ORANGE,
    )
    notes.append(SlideNotes("Finding 4: Cross-Validation Supports The Probe Signal", ["This slide adds confidence: layer selection is not based on a single noisy split."]))

    # 12
    slide = new_slide("Bridge: From Reading A Signal To Testing A Mechanism", "A probe can be accurate without being causal, so the next experiments intervene.")
    tests = [
        ("Steering", "Add a temporal direction and measure short-vs-long logit shift", BLUE),
        ("Activation patching", "Patch clean activations into corrupted prompts and measure recovery", TEAL),
        ("Attribution patching", "Use activation gradients as a first-order patch estimate", ORANGE),
        ("Ablation", "Remove or replace component outputs and measure drop", RED),
    ]
    for idx, (name, desc, color) in enumerate(tests):
        left = 0.8 + (idx % 2) * 6.1
        top = 1.45 + (idx // 2) * 2.15
        add_card(slide, left, top, 5.55, 1.55)
        add_text(slide, name, left + 0.25, top + 0.2, 5.0, 0.3, size=17, bold=True, color=color)
        add_text(slide, desc, left + 0.25, top + 0.68, 5.0, 0.5, size=12, color=INK)
    add_text(slide, "Interpretation rule: causal evidence is strongest when multiple intervention views point to the same model/layer/component pattern.", 1.1, 6.15, 11.2, 0.35, size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    notes.append(SlideNotes("Bridge: From Reading A Signal To Testing A Mechanism", ["Clarify why causal experiments are necessary after strong probe validation."]))

    # 13
    slide = new_slide("Causal Finding: Steering Moves The Logits", "Adding the temporal direction tests whether the representation is behaviorally active.")
    add_image(slide, RQ47_FIG_DIR / "steering_best_resid_heatmap.png", 0.75, 1.25, 7.1, 3.85)
    add_explanation(
        slide,
        "What The Heatmap Means",
        "Cells show the average logit shift from +3 steering versus no steering at the best residual layer. Larger positive values mean the temporal direction moves the model toward the targeted long-vs-short preference.",
        8.35,
        1.35,
        3.9,
        1.95,
        color=BLUE,
    )
    add_explanation(
        slide,
        "Why It Matters",
        "Steering is a direct causal manipulation. If adding the direction changes logits, the probe direction is not only readable; it is connected to the model's decision surface.",
        8.35,
        3.75,
        3.9,
        1.75,
        color=ORANGE,
    )
    notes.append(SlideNotes("Causal Finding: Steering Moves The Logits", ["Explain steering as a dose-response causal test, not merely another classifier metric."]))

    # 14
    slide = new_slide("Causal Finding: Activation Patching Recovers Behavior", "Clean activations can restore clean temporal choices on corrupted prompts.")
    add_image(slide, RQ47_FIG_DIR / "activation_patching_best_resid_heatmap.png", 0.75, 1.25, 7.1, 3.85)
    add_explanation(
        slide,
        "What The Heatmap Means",
        "Cells report median normalized recovery from patching clean residual activations into corrupted prompts. Values near 1 mean the patch recovers most of the clean-vs-corrupt logit difference.",
        8.35,
        1.35,
        3.9,
        1.95,
        color=TEAL,
    )
    add_explanation(
        slide,
        "Why It Matters",
        "Patching localizes function. When clean activations repair corrupted behavior, that layer/component is carrying information used by the downstream computation.",
        8.35,
        3.75,
        3.9,
        1.75,
        color=GREEN,
    )
    notes.append(SlideNotes("Causal Finding: Activation Patching Recovers Behavior", ["Use recovery as the intuitive metric: clean internal state repairs corrupted output tendency."]))

    # 15
    slide = new_slide("Causal Finding: Attribution Patching Is A Fast Cross-Check", "Gradient estimates often highlight similar residual-stream effects.")
    add_image(slide, RQ47_FIG_DIR / "attribution_patching_best_resid_heatmap.png", 0.75, 1.25, 7.1, 3.85)
    add_explanation(
        slide,
        "What The Heatmap Means",
        "Cells estimate the activation-patching effect using activation gradients. It is a first-order approximation rather than a full forward patched run.",
        8.35,
        1.35,
        3.9,
        1.75,
        color=ORANGE,
    )
    add_explanation(
        slide,
        "Why It Matters",
        "Agreement between activation patching and attribution patching is evidence that the identified layers are not a one-off artifact of one intervention method.",
        8.35,
        3.55,
        3.9,
        1.95,
        color=BLUE,
    )
    notes.append(SlideNotes("Causal Finding: Attribution Patching Is A Fast Cross-Check", ["Frame attribution patching as a triangulation tool, not as stronger than activation patching."]))

    # 16
    slide = new_slide("Causal Finding: Ablation Shows Dependence", "Removing residual information can reduce clean temporal behavior, with noisy edge cases.")
    add_image(slide, RQ47_FIG_DIR / "ablation_best_resid_heatmap.png", 0.75, 1.25, 7.1, 3.85)
    add_explanation(
        slide,
        "What The Heatmap Means",
        "Cells show median drop from clean behavior after ablating the selected residual layer. Positive values mean removing the component damaged the clean temporal signal.",
        8.35,
        1.35,
        3.9,
        1.75,
        color=RED,
    )
    add_explanation(
        slide,
        "Why It Matters",
        "Ablation asks the inverse question from patching: not can we add the signal back, but does removing it hurt? Together they make the causal story more complete.",
        8.35,
        3.55,
        3.9,
        1.95,
        color=ORANGE,
    )
    notes.append(SlideNotes("Causal Finding: Ablation Shows Dependence", ["Mention the denominator caveat verbally; ablation is useful but can have outliers."]))

    # 17
    slide = new_slide("Causal Synthesis", "The result is not one universal layer; it is a repeated residual-stream pattern with model-specific strength.")
    comp = bits["best_activation"][["model_alias", "layer_source_method", "median_effect"]].copy()
    comp["median_effect"] = comp["median_effect"].map(lambda v: f"{v:.2f}")
    comp = comp.rename(columns={"model_alias": "Model", "layer_source_method": "Selector", "median_effect": "Activation recovery"})
    add_simple_table(slide, comp, 0.85, 1.35, 5.9, 2.25, font_size=9)
    add_bullets(
        slide,
        [
            "Residual stream is the most consistent component-level site across steering, patching, and ablation.",
            "Model and selector matter: Phi-3 LR, Llama Attn-selected, and Qwen DMM/Attn patterns are the clearest cases.",
            "Top-k layer runs are best interpreted as robustness checks rather than the main headline.",
        ],
        7.15,
        1.45,
        5.1,
        2.45,
        size=14,
    )
    add_explanation(
        slide,
        "Why This Is The Right Claim",
        "The causal evidence is meaningful but heterogeneous. The careful conclusion is that temporal horizon is causally relevant in several settings, not that every model uses an identical circuit.",
        1.05,
        4.55,
        10.95,
        1.2,
        color=NAVY,
    )
    notes.append(SlideNotes("Causal Synthesis", ["This slide makes the causal story coherent and conservative."]))

    # 18
    slide = new_slide("Online Oversight Setup", "Probe scores are recorded at the prompt end and after every generated token.")
    xs = [0.9, 3.3, 5.7, 8.1, 10.5]
    labels = ["Prompt", "Generate token", "Extract hidden state", "Score probe", "Detect event"]
    for idx, (x, label) in enumerate(zip(xs, labels)):
        add_flow_box(slide, label, x, 2.0, 1.8, [BLUE, TEAL, ORANGE, GREEN, RED][idx])
        if idx < 4:
            add_arrow(slide, x + 1.88, 2.28)
    add_explanation(
        slide,
        "What The Monitor Measures",
        "For each prefix, it records the probe score at selected layers and marks the first step where a lightweight event detector fires.",
        1.0,
        3.65,
        5.6,
        1.1,
        color=BLUE,
    )
    add_explanation(
        slide,
        "Why It Matters",
        "This turns a static representation into an online signal. If score drift precedes the event, the representation may provide early warning rather than post-hoc explanation.",
        6.9,
        3.65,
        5.25,
        1.1,
        color=GREEN,
    )
    notes.append(SlideNotes("Online Oversight Setup", ["Describe this as an unfolding-process experiment: probe scores over generation time."]))

    # 19
    slide = new_slide("Online Finding: Pre-Event Drift Appears", "The heatmap summarizes median score increases before keyword-detected events.")
    add_image(slide, RQ47_FIG_DIR / "temporal_oversight_pre_event_delta_heatmap.png", 0.75, 1.25, 7.2, 3.9)
    add_explanation(
        slide,
        "What The Heatmap Means",
        "Each cell is the median increase from initial probe score to the maximum score before an event. Larger values mean stronger precursor drift.",
        8.35,
        1.35,
        3.9,
        1.55,
        color=GREEN,
    )
    add_explanation(
        slide,
        "Why It Matters",
        "The strongest drifts show that the temporal signal can move before the detector fires. This is evidence for a monitoring pipeline, not yet a final safety detector.",
        8.35,
        3.35,
        3.9,
        1.9,
        color=ORANGE,
    )
    notes.append(SlideNotes("Online Finding: Pre-Event Drift Appears", ["Keep the distinction clear: promising precursor signal, but keyword detector means no final safety claim."]))

    # 20
    slide = new_slide("Online Finding Detail: Top Pre-Event Drifts", "Concrete examples make the heatmap easier to interpret.")
    drift = bits["top_drift"][["model_alias", "prompt_id", "layer_source_method", "delta_before_event", "first_event_step"]]
    drift = drift.rename(columns={"model_alias": "Model", "prompt_id": "Prompt", "layer_source_method": "Sel.", "delta_before_event": "Drift", "first_event_step": "Event"})
    add_simple_table(slide, drift, 0.65, 1.25, 7.55, 3.25, font_size=8)
    add_explanation(
        slide,
        "What The Rows Mean",
        "Each row is one model/prompt/selector case where the probe score rose before the first keyword event. Drift is the pre-event score increase.",
        8.55,
        1.35,
        3.75,
        1.65,
        color=BLUE,
    )
    add_explanation(
        slide,
        "Why It Matters",
        "Examples like GPT-2 sycophancy and Phi-3 short-term pressure show the monitor can surface specific precursor traces for deeper qualitative review.",
        8.55,
        3.45,
        3.75,
        1.75,
        color=TEAL,
    )
    notes.append(SlideNotes("Online Finding Detail: Top Pre-Event Drifts", ["Use examples to make the online monitor tangible."]))

    # 21
    slide = new_slide("Limitations", "The evidence is useful, but the final safety claim needs stronger detectors and broader sampling.")
    add_bullets(
        slide,
        [
            "Keyword event detector validates the pipeline, not a final safety classifier.",
            "Online monitoring currently covers a small prompt set; more prompts are needed for confidence intervals.",
            "Ablation is component-level across Hugging Face models; per-head ablation remains future work.",
            "Some normalized patching and ablation scores are unstable when baseline denominators are small.",
            "The trajectory monitor recomputes full prefixes for clarity and portability; a KV-cache version would be faster.",
        ],
        0.9,
        1.45,
        11.6,
        4.25,
        size=17,
    )
    notes.append(SlideNotes("Limitations", ["This slide keeps the story honest and helps the technical audience trust the positive results."]))

    # 22
    slide = new_slide("Conclusion And Next Work", "Temporal horizon is decodable, causally relevant in several settings, and monitorable during generation.")
    add_bullets(
        slide,
        [
            "Decodability: hidden-state probes validate across GPT-2, Qwen3-4B, Phi-3-mini, and Llama-3.2-3B.",
            "Causality: steering, patching, attribution, and ablation identify meaningful residual-stream intervention sites.",
            "Online monitoring: pre-event probe drift appears in several trajectories.",
            "Next: replace keyword detector, join classifier scores, expand prompts, add per-head ablations, and quantify uncertainty.",
        ],
        0.9,
        1.45,
        8.35,
        3.45,
        size=17,
    )
    add_card(slide, 9.6, 1.65, 2.7, 2.8, fill=RGBColor(234, 243, 252))
    add_text(slide, "Bottom line", 9.9, 2.0, 2.1, 0.3, size=17, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(slide, "The signal is readable, often causally relevant, and usable as a monitoring trace.", 9.9, 2.55, 2.1, 1.15, size=14, color=INK, align=PP_ALIGN.CENTER)
    notes.append(SlideNotes("Conclusion And Next Work", ["End with the crisp claim: useful mechanistic oversight scaffold, with clear next experiments."]))

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(PPTX_PATH)
    return len(prs.slides), count_embedded_images(prs), notes


def count_embedded_images(prs: Presentation) -> int:
    total = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "shape_type", None) == 13:  # MSO_SHAPE_TYPE.PICTURE
                total += 1
    return total


def write_notes(notes: list[SlideNotes], slide_count: int, image_count: int) -> None:
    lines = [
        "# RQ47 Temporal Awareness Final Presentation Notes",
        "",
        f"Generated deck: `{PPTX_PATH.relative_to(ROOT)}`",
        f"Slide count: {slide_count}",
        f"Embedded image count: {image_count}",
        "",
        "## Speaker Notes",
        "",
    ]
    for idx, note in enumerate(notes, start=1):
        lines.append(f"### Slide {idx}: {note.title}")
        lines.extend(f"- {bullet}" for bullet in note.bullets)
        lines.append("")
    lines.extend(
        [
            "## Artifact Sources",
            "",
            "- Probe validation figures: `results/figures/probe_validation_multimethod/`",
            "- All-method validation figures: `results/figures/probe_validation_all_methods/`",
            "- RQ47 intervention and oversight figures: `results/figures/rq47/`",
            "- RQ47 summary tables: `results/tables/rq47/`",
            "",
            "## Verification Checklist",
            "",
            "- Generated PPTX is non-empty.",
            "- Slide count matches the expanded story structure.",
            "- Embedded image count is at least 12.",
            "- Notes markdown exists and records slide-level speaker guidance.",
        ]
    )
    NOTES_PATH.write_text("\n".join(lines), encoding="utf-8")


def main() -> None:
    slide_count, image_count, notes = build_deck()
    write_notes(notes, slide_count, image_count)
    print(f"Wrote {PPTX_PATH.relative_to(ROOT)}")
    print(f"Wrote {NOTES_PATH.relative_to(ROOT)}")
    print(f"Slides: {slide_count}")
    print(f"Embedded images: {image_count}")


if __name__ == "__main__":
    main()
